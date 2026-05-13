#!/usr/bin/env python3
"""
Generate synthetic PDF, DOCX, and PPTX files for the organize-messy-files training variant.

Categories (4 instead of 5):
1. machine_learning - 18 PDFs + 2 pptx
2. climate_change - 15 PDFs + 1 docx
3. neuroscience - 16 PDFs + 2 docx
4. ancient_history - 14 PDFs + 1 pptx + 1 xlsx

Total: 63 PDFs + 2 pptx + 3 docx + 1 xlsx = 69 files
"""

import os
import random
import string

from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from docx import Document
from pptx import Presentation
from pptx.util import Inches
import openpyxl

OUTPUT_DIR = "/home/zhanggenrui/workplace/self-evolving-skills/Benchmarks/skillsbench/tasks-train/organize-messy-files/environment/data"

# Content templates per category
ML_TOPICS = [
    "This paper presents a novel deep learning architecture for image classification using convolutional neural networks. We propose a residual attention mechanism that achieves state-of-the-art results on ImageNet benchmarks. The model uses gradient descent optimization with adaptive learning rates.",
    "We introduce a new approach to natural language processing using transformer models with self-attention mechanisms. Our method improves upon BERT and GPT architectures by incorporating sparse attention patterns for efficient training on large text corpora.",
    "This work explores reinforcement learning algorithms for robotic manipulation tasks. We combine policy gradient methods with model-based planning to achieve sample-efficient learning in continuous action spaces.",
    "A comprehensive study on generative adversarial networks for synthetic data generation. We propose a Wasserstein GAN variant with spectral normalization that produces high-fidelity images while maintaining training stability.",
    "We present a federated learning framework for privacy-preserving machine learning across distributed medical datasets. Our approach uses differential privacy guarantees while maintaining model accuracy.",
    "This paper introduces a graph neural network architecture for molecular property prediction. We leverage message passing neural networks with attention-based aggregation for drug discovery applications.",
    "An investigation into transfer learning techniques for low-resource language understanding. We demonstrate that multilingual pre-training significantly improves zero-shot cross-lingual classification performance.",
    "We propose a neural architecture search method using evolutionary algorithms. Our approach discovers efficient network topologies that outperform hand-designed architectures on CIFAR-10 and CIFAR-100.",
    "This study examines the interpretability of deep neural networks through gradient-based attribution methods. We develop a unified framework for explaining model predictions in computer vision tasks.",
    "A novel approach to semi-supervised learning using consistency regularization and pseudo-labeling. Our method achieves strong performance with only 10% labeled data on standard benchmarks.",
    "We introduce a variational autoencoder framework for anomaly detection in time series data. The model learns latent representations that capture normal patterns and flags deviations effectively.",
    "This paper presents advances in meta-learning for few-shot classification. We propose a task-agnostic feature extractor that generalizes across diverse visual recognition domains.",
    "An exploration of knowledge distillation techniques for model compression. We show that intermediate layer matching enables smaller student models to retain 95% of teacher accuracy.",
    "We develop a multi-task learning framework for simultaneous object detection and semantic segmentation. Our shared backbone architecture reduces computational costs while improving both tasks.",
    "This work proposes a contrastive learning approach for self-supervised visual representation learning. We demonstrate that momentum-based feature banks improve downstream task performance.",
    "A study on batch normalization alternatives for training deep residual networks. We introduce group normalization variants that perform consistently across different batch sizes.",
    "We present a Bayesian optimization framework for hyperparameter tuning of machine learning models. Our acquisition function balances exploration and exploitation efficiently in high-dimensional spaces.",
    "This paper investigates catastrophic forgetting in continual learning scenarios. We propose an elastic weight consolidation method that preserves important parameters while learning new tasks.",
]

CLIMATE_TOPICS = [
    "This study analyzes global temperature anomalies over the past century using satellite and ground-based measurements. We find accelerating warming trends in Arctic regions consistent with ice-albedo feedback amplification.",
    "We present a comprehensive analysis of greenhouse gas emissions from agricultural practices. Our life cycle assessment reveals that livestock production accounts for a significant portion of methane and nitrous oxide emissions.",
    "This paper models sea level rise projections under various representative concentration pathway scenarios. We incorporate thermal expansion and ice sheet dynamics to predict coastal flooding risks through 2100.",
    "An investigation into the effects of ocean acidification on coral reef ecosystems. We document declining calcification rates and biodiversity loss in tropical reef systems under elevated CO2 conditions.",
    "We develop a high-resolution regional climate model for predicting precipitation patterns in Southeast Asia. Our downscaling approach captures monsoon variability and extreme rainfall events more accurately.",
    "This study examines the carbon sequestration potential of reforestation programs in tropical regions. We quantify above-ground biomass accumulation and soil carbon changes over 20-year restoration timelines.",
    "A comprehensive review of renewable energy technologies and their role in climate change mitigation. We assess the levelized cost of energy for solar, wind, and geothermal sources across different regions.",
    "We analyze the impact of permafrost thawing on methane release in Siberian tundra. Our field measurements indicate accelerating organic matter decomposition and positive feedback to global warming.",
    "This paper presents new paleoclimate data from Antarctic ice cores spanning 800,000 years. We reconstruct atmospheric CO2 concentrations and temperature variations to validate current climate models.",
    "An assessment of urban heat island effects on energy consumption and public health in megacities. We propose green infrastructure interventions to reduce ambient temperatures by 2-4 degrees Celsius.",
    "We investigate the relationship between deforestation in the Amazon basin and regional rainfall patterns. Our analysis shows that large-scale forest loss reduces precipitation recycling and extends dry seasons.",
    "This study quantifies the global carbon budget and remaining emissions allowance for 1.5 degree Celsius warming targets. We incorporate updated climate sensitivity estimates and land use change data.",
    "A modeling study of thermohaline circulation weakening under future warming scenarios. We assess the implications for European climate and Atlantic fisheries productivity.",
    "We present satellite-based observations of glacier retreat in the Himalayan region over the past four decades. Our mass balance calculations indicate accelerating ice loss threatening water resources for billions.",
    "This paper evaluates the effectiveness of carbon capture and storage technologies for mitigating industrial emissions. We compare geological sequestration, direct air capture, and enhanced weathering approaches.",
]

NEURO_TOPICS = [
    "This study investigates the role of hippocampal place cells in spatial navigation and memory formation. We use optogenetic manipulation to demonstrate causal relationships between neural activity patterns and behavioral performance.",
    "We present a comprehensive connectome analysis of the Drosophila visual system. Our electron microscopy reconstruction reveals novel synaptic motifs in motion detection circuits.",
    "This paper examines the molecular mechanisms of long-term potentiation at glutamatergic synapses. We identify a novel phosphorylation cascade involving CaMKII that modulates AMPA receptor trafficking.",
    "An investigation into the neural correlates of consciousness using intracranial EEG recordings in epilepsy patients. We identify prefrontal-parietal gamma oscillations that correlate with subjective awareness.",
    "We develop a brain-computer interface using high-density microelectrode arrays for decoding motor intentions. Our algorithm achieves 95% accuracy in predicting hand movements from primary motor cortex signals.",
    "This study explores the role of astrocytes in synaptic transmission and neural circuit function. We demonstrate that calcium signaling in astrocytic networks modulates neurotransmitter release probability.",
    "A neuroimaging study of default mode network alterations in major depressive disorder. We identify disrupted functional connectivity patterns that predict treatment response to selective serotonin reuptake inhibitors.",
    "We present a computational model of cerebellar learning based on Purkinje cell dendritic computation. Our simulation reproduces adaptive timing in classical conditioning paradigms.",
    "This paper investigates adult neurogenesis in the subventricular zone and its contribution to olfactory discrimination. We trace newborn neurons using retroviral labeling and assess their functional integration.",
    "An analysis of sleep spindle dynamics and their role in memory consolidation. We show that thalamocortical oscillations during NREM sleep facilitate hippocampal-neocortical dialogue.",
    "We examine the neuropharmacology of dopamine signaling in reward prediction and decision making. Our microdialysis measurements reveal phasic dopamine release patterns in the nucleus accumbens.",
    "This study maps the somatotopic organization of the primary somatosensory cortex using high-field functional MRI. We reveal fine-grained digit representations that reorganize following peripheral nerve injury.",
    "A systematic review of neuroinflammation in Alzheimer's disease pathogenesis. We evaluate the roles of microglia activation, complement system, and inflammatory cytokines in amyloid-beta clearance.",
    "We develop a novel neural probe technology for chronic multi-region recordings in freely moving rodents. Our flexible polymer electrode arrays minimize tissue damage and maintain signal quality for months.",
    "This paper presents evidence for predictive coding mechanisms in the visual cortex. We use computational modeling and electrophysiology to show that feedback connections carry prediction error signals.",
    "An investigation into the genetic basis of autism spectrum disorders using genome-wide association studies. We identify novel risk loci affecting synaptic adhesion molecules and chromatin remodeling factors.",
]

HISTORY_TOPICS = [
    "This paper examines the administrative structures of the Roman Empire during the Principate era. We analyze epigraphic evidence from provincial inscriptions to reconstruct bureaucratic hierarchies and taxation systems.",
    "A comprehensive study of ancient Egyptian mummification techniques using CT scanning and chemical analysis. We identify previously unknown embalming substances and trace trade networks for aromatic resins.",
    "We investigate the collapse of the Bronze Age civilizations in the Eastern Mediterranean circa 1200 BCE. Our analysis integrates archaeological, climatological, and textual evidence to evaluate multiple causation theories.",
    "This study presents new archaeological findings from Mesopotamian cuneiform tablets documenting mathematical knowledge in ancient Babylon. We translate algebraic problem texts that predate Greek mathematics by centuries.",
    "An analysis of trade routes connecting Han Dynasty China with the Roman Empire along the Silk Road. We examine numismatic evidence and material culture to reconstruct patterns of long-distance exchange.",
    "We reexamine the fall of the Western Roman Empire through the lens of environmental history. Our pollen analysis and dendrochronology data reveal agricultural disruptions preceding barbarian incursions.",
    "This paper documents newly discovered cave paintings in the Dordogne region of France dating to the Upper Paleolithic. We use uranium-thorium dating to establish a chronology spanning 35,000 to 15,000 years ago.",
    "A study of ancient Greek democratic institutions and their evolution from the reforms of Cleisthenes to the Hellenistic period. We analyze voting records and assembly proceedings preserved in Athenian inscriptions.",
    "We investigate the construction techniques of Angkor Wat and surrounding Khmer temples using LIDAR survey data. Our analysis reveals a vast urban infrastructure hidden beneath tropical vegetation.",
    "This study examines the origins of writing systems in ancient Sumer through analysis of proto-cuneiform tokens. We trace the transition from accounting devices to phonetic script over two millennia.",
    "An archaeological investigation of Viking settlement patterns in Iceland and Greenland. We combine radiocarbon dating with saga literature to reconstruct colonization timelines and subsistence strategies.",
    "We present a comprehensive analysis of Maya hieroglyphic texts from Palenque documenting royal succession and warfare. Our new decipherments reveal previously unknown political alliances between Classic Maya city-states.",
    "This paper examines the social and economic impact of the Black Death on medieval European societies. We analyze demographic records, price indices, and labor legislation to assess long-term structural changes.",
    "A study of ancient Indus Valley Civilization urban planning using geophysical survey methods. We map the street grid and drainage systems of Mohenjo-daro revealing sophisticated civil engineering.",
]

# File naming: use randomized document IDs (different from arxiv style)
def gen_id():
    """Generate a random document ID like DOC-A3X7K2."""
    chars = string.ascii_uppercase + string.digits
    return "DOC-" + ''.join(random.choice(chars) for _ in range(6))

random.seed(42)

# Generate unique IDs for all PDFs
def gen_unique_ids(n):
    ids = set()
    while len(ids) < n:
        ids.add(gen_id())
    return sorted(ids)

ML_IDS = gen_unique_ids(18)
CLIMATE_IDS = gen_unique_ids(15)
NEURO_IDS = gen_unique_ids(16)
HISTORY_IDS = gen_unique_ids(14)

# Make sure no collisions across categories
all_ids = ML_IDS + CLIMATE_IDS + NEURO_IDS + HISTORY_IDS
assert len(all_ids) == len(set(all_ids)), "ID collision!"

def create_pdf(filepath, title, content):
    c = canvas.Canvas(filepath, pagesize=letter)
    c.setFont("Helvetica-Bold", 16)
    c.drawString(72, 720, title)
    c.setFont("Helvetica", 11)

    # Write content with word wrapping
    y = 680
    words = content.split()
    line = ""
    for word in words:
        test_line = f"{line} {word}".strip()
        if len(test_line) > 90:
            c.drawString(72, y, line)
            y -= 15
            line = word
            if y < 72:
                c.showPage()
                c.setFont("Helvetica", 11)
                y = 720
        else:
            line = test_line
    if line:
        c.drawString(72, y, line)

    c.save()

def create_docx(filepath, title, content):
    doc = Document()
    doc.add_heading(title, 0)
    doc.add_paragraph(content)
    doc.save(filepath)

def create_pptx(filepath, title, content):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # Title and content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content
    prs.save(filepath)

def create_xlsx(filepath, title, data_rows):
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = title
    for row in data_rows:
        ws.append(row)
    wb.save(filepath)

def main():
    os.makedirs(OUTPUT_DIR, exist_ok=True)

    # Generate ML PDFs
    for i, doc_id in enumerate(ML_IDS):
        filepath = os.path.join(OUTPUT_DIR, f"{doc_id}.pdf")
        create_pdf(filepath, f"Machine Learning Research Paper {i+1}", ML_TOPICS[i])

    # Generate Climate PDFs
    for i, doc_id in enumerate(CLIMATE_IDS):
        filepath = os.path.join(OUTPUT_DIR, f"{doc_id}.pdf")
        create_pdf(filepath, f"Climate Change Research Paper {i+1}", CLIMATE_TOPICS[i])

    # Generate Neuroscience PDFs
    for i, doc_id in enumerate(NEURO_IDS):
        filepath = os.path.join(OUTPUT_DIR, f"{doc_id}.pdf")
        create_pdf(filepath, f"Neuroscience Research Paper {i+1}", NEURO_TOPICS[i])

    # Generate History PDFs
    for i, doc_id in enumerate(HISTORY_IDS):
        filepath = os.path.join(OUTPUT_DIR, f"{doc_id}.pdf")
        create_pdf(filepath, f"Ancient History Research Paper {i+1}", HISTORY_TOPICS[i])

    # Generate PPTX files for machine_learning
    create_pptx(
        os.path.join(OUTPUT_DIR, "NeurIPS_talk.pptx"),
        "Deep Learning Advances at NeurIPS",
        "Recent breakthroughs in neural network architectures, training methodology, and machine learning applications for computer vision and NLP."
    )
    create_pptx(
        os.path.join(OUTPUT_DIR, "ICML_workshop.pptx"),
        "Reinforcement Learning Workshop at ICML",
        "Policy optimization, multi-agent reinforcement learning, and applications of deep RL to robotics and game playing."
    )

    # Generate DOCX files for climate_change
    create_docx(
        os.path.join(OUTPUT_DIR, "ipcc_summary_draft.docx"),
        "IPCC Working Group Summary Draft",
        "This document summarizes the latest findings on anthropogenic greenhouse gas emissions, global mean temperature projections, and recommended climate change mitigation strategies. Sea level rise estimates are updated based on new ice sheet modeling."
    )

    # Generate DOCX files for neuroscience
    create_docx(
        os.path.join(OUTPUT_DIR, "lab_protocol_ephys.docx"),
        "Electrophysiology Lab Protocol",
        "Standard operating procedure for in vivo extracellular recordings in rodent hippocampus. Includes electrode preparation, stereotaxic surgery coordinates, signal amplification settings, and spike sorting methodology for neural data analysis."
    )
    create_docx(
        os.path.join(OUTPUT_DIR, "grant_proposal_brain.docx"),
        "NIH R01 Grant Proposal: Neural Circuit Mapping",
        "This proposal outlines a research program to map functional connectivity in the prefrontal cortex using two-photon calcium imaging. We will investigate how neural ensembles encode working memory and decision-making in behaving mice."
    )

    # Generate PPTX for ancient_history
    create_pptx(
        os.path.join(OUTPUT_DIR, "archaeology_lecture.pptx"),
        "Introduction to Classical Archaeology",
        "Survey of archaeological methods applied to ancient Mediterranean civilizations. Covers excavation techniques, stratigraphy, pottery typology, and dating methods for Roman and Greek sites."
    )

    # Generate XLSX for ancient_history
    create_xlsx(
        os.path.join(OUTPUT_DIR, "artifact_catalog.xlsx"),
        "Archaeological Artifacts",
        [
            ["Artifact ID", "Site", "Period", "Type", "Material", "Date Found"],
            ["ART-001", "Pompeii", "Roman Imperial", "Amphora", "Ceramic", "2019-03-15"],
            ["ART-002", "Knossos", "Minoan", "Fresco Fragment", "Plaster", "2018-07-22"],
            ["ART-003", "Troy", "Bronze Age", "Spearhead", "Bronze", "2020-01-10"],
            ["ART-004", "Athens", "Classical Greek", "Red-figure Kylix", "Ceramic", "2017-09-05"],
            ["ART-005", "Carthage", "Punic", "Stele", "Limestone", "2021-04-18"],
            ["ART-006", "Ephesus", "Hellenistic", "Coin Hoard", "Silver", "2019-11-30"],
            ["ART-007", "Luxor", "New Kingdom", "Scarab Seal", "Faience", "2016-06-12"],
        ]
    )

    total = len(os.listdir(OUTPUT_DIR))
    print(f"Generated {total} files in {OUTPUT_DIR}")

    # Print the mapping for reference
    print("\n--- FILE MAPPING ---")
    print(f"machine_learning PDFs: {ML_IDS}")
    print(f"climate_change PDFs: {CLIMATE_IDS}")
    print(f"neuroscience PDFs: {NEURO_IDS}")
    print(f"ancient_history PDFs: {HISTORY_IDS}")

if __name__ == "__main__":
    main()
